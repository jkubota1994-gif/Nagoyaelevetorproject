import os
import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_qr_code(url, filename):
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(filename)

def add_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Bullet point layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

def generate_pptx():
    from pptx.util import Cm
    prs = Presentation()
    
    # Set slide size to A4 (21.0cm x 29.7cm)
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Header (Title)
    shape_title = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(19), Cm(2))
    tf_title = shape_title.text_frame
    p_title = tf_title.add_paragraph()
    p_title.text = "名古屋市営地下鉄エレベーターマップ"
    p_title.font.bold = True
    p_title.font.size = Pt(32)
    p_title.alignment = PP_ALIGN.CENTER

    p_sub = tf_title.add_paragraph()
    p_sub.text = "〜 お出かけをもっと自由に、もっとスムーズに 〜"
    p_sub.font.size = Pt(16)
    p_sub.alignment = PP_ALIGN.CENTER

    # 2. Mockup Image (中央上部)
    mockup_path = r"C:\Users\break\.gemini\antigravity\brain\6ab1cf61-f21a-483e-9d64-9d91c8dc1bd7\elevator_map_app_mockup_1774268923599.png"
    if os.path.exists(mockup_path):
        slide.shapes.add_picture(mockup_path, Cm(4.5), Cm(3), Cm(12), Cm(8))

    # 3. Section: Detailed Usage (中央〜下部)
    shape_usage = slide.shapes.add_textbox(Cm(1), Cm(11.5), Cm(19), Cm(10))
    tf_usage = shape_usage.text_frame
    tf_usage.word_wrap = True
    
    p_h = tf_usage.add_paragraph()
    p_h.text = "【 詳しい使い方 】"
    p_h.font.bold = True
    p_h.font.size = Pt(22)
    p_h.alignment = PP_ALIGN.CENTER

    steps = [
        ("① 駅名や目的地で検索", "画面上部の検索バーに目的地や駅名を入力。候補から選ぶだけで、地図がその場所にジャンプします。"),
        ("② エレベーター情報をチェック", "地図上の「青いピン」をタップ！そのエレベーターが「地上⇔改札」か「改札⇔ホーム」かなど、詳細な接続情報を確認できます。"),
        ("③ バリアフリールートを検索", "「ルート検索」ボタンから目的地を設定。階段を使わない、エレベーター優先の安心な経路を自動で計算して表示します。"),
        ("④ 駅パネルで次発時刻を確認", "駅名をタップすると表示されるパネルから、次に発車する列車の時刻をリアルタイムでチェック。待ち時間を最小限に。")
    ]
    
    for title, desc in steps:
        p_title = tf_usage.add_paragraph()
        p_title.text = title
        p_title.font.bold = True
        p_title.font.size = Pt(16)
        p_title.space_before = Pt(10)
        
        p_desc = tf_usage.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.space_after = Pt(5)

    # 4. Section: Access & QR Code (最下部)
    url = "https://jkubota1994-gif.github.io/Nagoyaelevetorproject/"
    qr_filename = "github_pages_qr.png"
    create_qr_code(url, qr_filename)

    # QR Code (右下)
    slide.shapes.add_picture(qr_filename, Cm(14), Cm(22.5), Cm(5.5), Cm(5.5))

    # URL & Info (左下)
    shape_footer = slide.shapes.add_textbox(Cm(1), Cm(23.5), Cm(12), Cm(4))
    tf_footer = shape_footer.text_frame
    p_footer = tf_footer.add_paragraph()
    p_footer.text = "今すぐアクセスして\n便利さを体験してください！"
    p_footer.font.bold = True
    p_footer.font.size = Pt(18)
    
    p_url = tf_footer.add_paragraph()
    p_url.text = f"\nURL: {url}"
    p_url.font.size = Pt(11)

    # Save pptx
    output_filename = "Nagoya_Elevator_Map_Pamphlet_Final.pptx"
    prs.save(output_filename)
    print(f"Detailed A4 Pamphlet generated: {output_filename}")
    
    # Clean up QR image
    if os.path.exists(qr_filename):
        os.remove(qr_filename)

if __name__ == "__main__":
    generate_pptx()
